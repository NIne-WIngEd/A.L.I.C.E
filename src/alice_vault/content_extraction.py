from __future__ import annotations

import csv
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable


SUPPORTED_FAMILIES = {
    "json", "html", "csv", "text", "pdf", "docx", "xlsx", "pptx",
    "calendar", "contacts", "subtitles", "xml",
}


@dataclass(frozen=True)
class ExtractionResult:
    status: str
    text: str
    chars: int
    truncated: bool
    parser: str
    error: str | None = None


def _cap(text: str, max_chars: int) -> tuple[str, bool]:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{4,}", "\n\n\n", text).strip()
    if len(text) <= max_chars:
        return text, False
    head = max_chars * 2 // 3
    tail = max_chars - head
    return text[:head] + "\n\n[...TRUNCATED...]\n\n" + text[-tail:], True


def _decode_bytes(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _iter_json_strings(value: Any, prefix: str = "") -> Iterable[str]:
    if isinstance(value, dict):
        for key, item in value.items():
            next_prefix = f"{prefix}.{key}" if prefix else str(key)
            yield from _iter_json_strings(item, next_prefix)
    elif isinstance(value, list):
        for index, item in enumerate(value):
            yield from _iter_json_strings(item, f"{prefix}[{index}]")
    elif value is None:
        return
    elif isinstance(value, (str, int, float, bool)):
        text = str(value).strip()
        if text:
            yield f"{prefix}: {text}" if prefix else text


def _extract_json(path: Path) -> str:
    data = json.loads(_decode_bytes(path.read_bytes()))
    return "\n".join(_iter_json_strings(data))


def _extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode_bytes(path.read_bytes()), "html.parser")
    for tag in soup(["script", "style", "noscript", "template"]):
        tag.decompose()
    return soup.get_text("\n")


def _extract_xml(path: Path) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode_bytes(path.read_bytes()), "xml")
    return soup.get_text("\n")


def _extract_csv(path: Path, max_rows: int = 2000, max_cols: int = 80) -> str:
    text = _decode_bytes(path.read_bytes())
    lines: list[str] = []
    reader = csv.reader(text.splitlines())
    for row_index, row in enumerate(reader):
        if row_index >= max_rows:
            lines.append("[ROW LIMIT REACHED]")
            break
        lines.append(" | ".join(cell.strip() for cell in row[:max_cols]))
    return "\n".join(lines)


def _extract_pdf(path: Path, max_pages: int = 40) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        try:
            if reader.decrypt("") == 0:
                raise ValueError("Encrypted PDF requires a password")
        except Exception as exc:
            raise ValueError("Encrypted PDF could not be opened") from exc
    parts: list[str] = []
    for index, page in enumerate(reader.pages[:max_pages]):
        parts.append(f"[PAGE {index + 1}]\n{page.extract_text() or ''}")
    if len(reader.pages) > max_pages:
        parts.append(f"[PAGE LIMIT REACHED: {len(reader.pages)} total pages]")
    return "\n".join(parts)


def _extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _extract_xlsx(
    path: Path,
    max_sheets: int = 20,
    max_rows_per_sheet: int = 500,
    max_cols: int = 80,
) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet_index, sheet in enumerate(workbook.worksheets):
            if sheet_index >= max_sheets:
                parts.append("[SHEET LIMIT REACHED]")
                break
            parts.append(f"[SHEET: {sheet.title}]")
            for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                if row_index >= max_rows_per_sheet:
                    parts.append("[ROW LIMIT REACHED]")
                    break
                values = ["" if value is None else str(value) for value in row[:max_cols]]
                if any(value.strip() for value in values):
                    parts.append(" | ".join(values))
    finally:
        workbook.close()
    return "\n".join(parts)


def _extract_pptx(path: Path, max_slides: int = 100) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide_index, slide in enumerate(presentation.slides):
        if slide_index >= max_slides:
            parts.append("[SLIDE LIMIT REACHED]")
            break
        parts.append(f"[SLIDE {slide_index + 1}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def extract_text(
    path: Path,
    family: str,
    *,
    max_chars: int = 24_000,
    max_file_bytes: int = 30 * 1024 * 1024,
) -> ExtractionResult:
    family = family.strip().lower()
    path = path.resolve(strict=True)
    if not path.is_file():
        return ExtractionResult("error", "", 0, False, family, "Not a file")
    size = path.stat().st_size
    if size == 0:
        return ExtractionResult("empty", "", 0, False, family)
    if size > max_file_bytes:
        return ExtractionResult(
            "too_large", "", 0, False, family,
            f"File exceeds {max_file_bytes} byte extraction limit",
        )
    if family not in SUPPORTED_FAMILIES:
        return ExtractionResult("unsupported", "", 0, False, family)

    try:
        if family == "json":
            text = _extract_json(path)
        elif family == "html":
            text = _extract_html(path)
        elif family == "xml":
            text = _extract_xml(path)
        elif family == "csv":
            text = _extract_csv(path)
        elif family == "pdf":
            text = _extract_pdf(path)
        elif family == "docx":
            text = _extract_docx(path)
        elif family == "xlsx":
            text = _extract_xlsx(path)
        elif family == "pptx":
            text = _extract_pptx(path)
        else:
            text = _decode_bytes(path.read_bytes())
        capped, truncated = _cap(text, max_chars)
        if not capped.strip():
            return ExtractionResult("no_text", "", 0, False, family)
        return ExtractionResult("ok", capped, len(capped), truncated, family)
    except Exception as exc:
        return ExtractionResult(
            "error", "", 0, False, family,
            f"{type(exc).__name__}: {exc}",
        )
