from __future__ import annotations

import csv
import json
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from .parser_registry import ParserSpec


@dataclass(frozen=True)
class ParsedDocument:
    text: str
    metadata: dict[str, Any]
    warnings: list[str]
    truncated: bool


@dataclass
class TextCollector:
    max_chars: int
    parts: list[str] = field(default_factory=list)
    chars: int = 0
    truncated: bool = False

    def add(self, value: Any) -> bool:
        if self.truncated:
            return False
        text = "" if value is None else str(value)
        text = text.replace("\x00", " ")
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        text = "".join(
            character
            for character in text
            if character in "\n\t" or ord(character) >= 32
        )
        text = re.sub(r"[ \t]+\n", "\n", text).strip()
        if not text:
            return True

        remaining = self.max_chars - self.chars
        separator = "\n" if self.parts else ""
        required = len(separator) + len(text)
        if required <= remaining:
            if separator:
                self.parts.append(separator)
            self.parts.append(text)
            self.chars += required
            return True

        if remaining > len(separator):
            if separator:
                self.parts.append(separator)
                remaining -= len(separator)
                self.chars += len(separator)
            self.parts.append(text[:remaining])
            self.chars += remaining
        self.truncated = True
        return False

    def build(self) -> str:
        return "".join(self.parts).strip()


def _decode_bytes(data: bytes) -> str:
    for encoding in (
        "utf-8-sig",
        "utf-8",
        "utf-16",
        "cp1252",
        "latin-1",
    ):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _validate_ooxml(path: Path, required_member: str) -> None:
    if not zipfile.is_zipfile(path):
        raise ValueError("Expected an Office Open XML ZIP container")
    with zipfile.ZipFile(path, "r") as archive:
        names = set(archive.namelist())
        if required_member not in names:
            raise ValueError(
                f"OOXML container is missing {required_member!r}"
            )


def _parse_json(
    path: Path,
    collector: TextCollector,
    spec: ParserSpec,
) -> dict[str, Any]:
    value = json.loads(_decode_bytes(path.read_bytes()))
    max_nodes = int(spec.limits.get("max_nodes", 150000))
    max_depth = int(spec.limits.get("max_depth", 80))
    stack: list[tuple[str, Any, int]] = [("", value, 0)]
    nodes = 0

    while stack and not collector.truncated:
        prefix, current, depth = stack.pop()
        nodes += 1
        if nodes > max_nodes:
            collector.truncated = True
            break
        if depth > max_depth:
            collector.add(f"{prefix}: [DEPTH LIMIT]")
            continue

        if isinstance(current, dict):
            for key, child in reversed(list(current.items())):
                next_prefix = f"{prefix}.{key}" if prefix else str(key)
                stack.append((next_prefix, child, depth + 1))
        elif isinstance(current, list):
            for index in range(len(current) - 1, -1, -1):
                stack.append(
                    (f"{prefix}[{index}]", current[index], depth + 1)
                )
        elif current is not None:
            collector.add(
                f"{prefix}: {current}" if prefix else str(current)
            )

    return {"nodes_seen": nodes}


def _parse_html(
    path: Path,
    collector: TextCollector,
) -> dict[str, Any]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode_bytes(path.read_bytes()), "html.parser")
    removed = 0
    for tag in soup(
        [
            "script",
            "style",
            "noscript",
            "template",
            "iframe",
            "object",
            "embed",
        ]
    ):
        tag.decompose()
        removed += 1
    collector.add(soup.get_text("\n"))
    return {"active_or_nontext_elements_removed": removed}


def _parse_csv(
    path: Path,
    collector: TextCollector,
    spec: ParserSpec,
) -> dict[str, Any]:
    max_rows = int(spec.limits.get("max_rows", 100000))
    max_columns = int(spec.limits.get("max_columns", 200))
    text = _decode_bytes(path.read_bytes())
    rows = 0
    for row in csv.reader(text.splitlines()):
        if rows >= max_rows:
            collector.truncated = True
            break
        collector.add(
            " | ".join(str(cell).strip() for cell in row[:max_columns])
        )
        rows += 1
        if collector.truncated:
            break
    return {"rows_extracted": rows}


def _parse_text(path: Path, collector: TextCollector) -> dict[str, Any]:
    collector.add(_decode_bytes(path.read_bytes()))
    return {}


def _parse_pdf(
    path: Path,
    collector: TextCollector,
    spec: ParserSpec,
) -> dict[str, Any]:
    from pypdf import PdfReader

    with path.open("rb") as handle:
        if handle.read(5) != b"%PDF-":
            raise ValueError("PDF signature is missing")

    reader = PdfReader(str(path), strict=False)
    if reader.is_encrypted:
        try:
            if reader.decrypt("") == 0:
                raise ValueError("Encrypted PDF requires a password")
        except Exception as exc:
            raise ValueError("Encrypted PDF could not be opened") from exc

    max_pages = int(spec.limits.get("max_pages", 200))
    max_page_content_bytes = int(
        spec.limits.get("max_page_content_bytes", 52428800)
    )
    pages_extracted = 0
    oversized_pages = 0

    for index, page in enumerate(reader.pages):
        if index >= max_pages:
            collector.truncated = True
            break

        contents = page.get_contents()
        if contents is not None:
            try:
                content_size = len(contents.get_data())
            except Exception:
                content_size = 0
            if content_size > max_page_content_bytes:
                oversized_pages += 1
                collector.add(
                    f"[PAGE {index + 1} SKIPPED: content stream exceeds limit]"
                )
                continue

        collector.add(f"[PAGE {index + 1}]")
        collector.add(page.extract_text() or "")
        pages_extracted += 1
        if collector.truncated:
            break

    return {
        "total_pages": len(reader.pages),
        "pages_extracted": pages_extracted,
        "oversized_pages_skipped": oversized_pages,
        "ocr_performed": False,
    }


def _parse_docx(
    path: Path,
    collector: TextCollector,
    spec: ParserSpec,
) -> dict[str, Any]:
    from docx import Document

    _validate_ooxml(path, "word/document.xml")
    document = Document(str(path))
    max_blocks = int(spec.limits.get("max_blocks", 100000))
    blocks = 0

    for paragraph in document.paragraphs:
        if blocks >= max_blocks or collector.truncated:
            collector.truncated = True
            break
        collector.add(paragraph.text)
        blocks += 1

    for table in document.tables:
        for row in table.rows:
            if blocks >= max_blocks or collector.truncated:
                collector.truncated = True
                break
            collector.add(" | ".join(cell.text.strip() for cell in row.cells))
            blocks += 1
        if collector.truncated:
            break

    return {
        "paragraphs": len(document.paragraphs),
        "top_level_tables": len(document.tables),
        "blocks_extracted": blocks,
    }


def _parse_xlsx(
    path: Path,
    collector: TextCollector,
    spec: ParserSpec,
) -> dict[str, Any]:
    from openpyxl import load_workbook

    _validate_ooxml(path, "xl/workbook.xml")
    workbook = load_workbook(
        str(path),
        read_only=True,
        data_only=True,
        keep_links=False,
    )
    max_sheets = int(spec.limits.get("max_sheets", 50))
    max_rows = int(spec.limits.get("max_rows_per_sheet", 10000))
    max_columns = int(spec.limits.get("max_columns", 200))
    sheets = 0
    rows = 0

    try:
        for sheet_index, sheet in enumerate(workbook.worksheets):
            if sheet_index >= max_sheets:
                collector.truncated = True
                break
            collector.add(f"[SHEET: {sheet.title}]")
            sheets += 1
            for row_index, row in enumerate(
                sheet.iter_rows(values_only=True)
            ):
                if row_index >= max_rows:
                    collector.truncated = True
                    break
                values = [
                    "" if value is None else str(value)
                    for value in row[:max_columns]
                ]
                if any(value.strip() for value in values):
                    collector.add(" | ".join(values))
                rows += 1
                if collector.truncated:
                    break
            if collector.truncated:
                break
    finally:
        workbook.close()

    return {
        "sheets_extracted": sheets,
        "rows_seen": rows,
        "read_only": True,
        "data_only": True,
        "external_links_loaded": False,
    }


def _parse_pptx(
    path: Path,
    collector: TextCollector,
    spec: ParserSpec,
) -> dict[str, Any]:
    from pptx import Presentation

    _validate_ooxml(path, "ppt/presentation.xml")
    presentation = Presentation(str(path))
    max_slides = int(spec.limits.get("max_slides", 300))
    slides = 0

    for index, slide in enumerate(presentation.slides):
        if index >= max_slides:
            collector.truncated = True
            break
        collector.add(f"[SLIDE {index + 1}]")
        slides += 1
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                collector.add(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    collector.add(
                        " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                    )
            if collector.truncated:
                break
        if collector.truncated:
            break

    return {
        "total_slides": len(presentation.slides),
        "slides_extracted": slides,
    }


def _unfold_lines(text: str) -> list[str]:
    output: list[str] = []
    for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        if line.startswith((" ", "\t")) and output:
            output[-1] += line[1:]
        else:
            output.append(line)
    return output


def _parse_contacts(
    path: Path,
    collector: TextCollector,
) -> dict[str, Any]:
    skipped_photos = 0
    records = 0
    for line in _unfold_lines(_decode_bytes(path.read_bytes())):
        upper = line.upper()
        if upper.startswith("BEGIN:VCARD"):
            records += 1
        if upper.startswith("PHOTO"):
            skipped_photos += 1
            continue
        collector.add(line)
        if collector.truncated:
            break
    return {
        "vcard_records": records,
        "photo_payloads_skipped": skipped_photos,
    }


def _parse_calendar(
    path: Path,
    collector: TextCollector,
) -> dict[str, Any]:
    skipped_attachments = 0
    events = 0
    for line in _unfold_lines(_decode_bytes(path.read_bytes())):
        upper = line.upper()
        if upper.startswith("BEGIN:VEVENT"):
            events += 1
        if upper.startswith("ATTACH"):
            skipped_attachments += 1
            continue
        collector.add(line)
        if collector.truncated:
            break
    return {
        "events": events,
        "attachments_skipped": skipped_attachments,
    }


def parse_document(path: Path, spec: ParserSpec) -> ParsedDocument:
    path = path.expanduser().resolve(strict=True)
    if not path.is_file():
        raise ValueError("Source is not a regular file")
    size = path.stat().st_size
    if size == 0:
        raise ValueError("Source file is empty")
    if size > spec.max_file_bytes:
        raise ValueError(
            f"Source exceeds parser limit of {spec.max_file_bytes} bytes"
        )

    collector = TextCollector(max_chars=spec.max_output_chars)
    family = spec.family

    if family == "json":
        metadata = _parse_json(path, collector, spec)
    elif family == "html":
        metadata = _parse_html(path, collector)
    elif family == "csv":
        metadata = _parse_csv(path, collector, spec)
    elif family in {"text", "subtitles"}:
        metadata = _parse_text(path, collector)
    elif family == "pdf":
        metadata = _parse_pdf(path, collector, spec)
    elif family == "docx":
        metadata = _parse_docx(path, collector, spec)
    elif family == "xlsx":
        metadata = _parse_xlsx(path, collector, spec)
    elif family == "pptx":
        metadata = _parse_pptx(path, collector, spec)
    elif family == "contacts":
        metadata = _parse_contacts(path, collector)
    elif family == "calendar":
        metadata = _parse_calendar(path, collector)
    else:
        raise ValueError(f"No implementation for parser family {family!r}")

    text = collector.build()
    if not text:
        raise ValueError("Parser produced no text")

    warnings: list[str] = []
    if collector.truncated:
        warnings.append("output_character_limit_reached")

    return ParsedDocument(
        text=text,
        metadata=metadata,
        warnings=warnings,
        truncated=collector.truncated,
    )
